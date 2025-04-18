import streamlit as st
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import io
import matplotlib.pyplot as plt
import base64
import pandas as pd
import plotly.express as px
import plotly.io as pio
from streamlit_sortables import sort_items

# Set default format for Plotly image export
pio.kaleido.scope.default_format = "png"

# Custom CSS for visual polish
st.markdown(
    """
    <style>
    .stButton>button {
        background-color: #4CAF50;
        color: white;
        border-radius: 8px;
        padding: 8px 16px;
        transition: background-color 0.3s;
    }
    .stButton>button:hover {
        background-color: #45a049;
    }
    .stTab {
        background-color: #f0f2f6;
        border-radius: 8px;
        padding: 10px;
    }
    .stMarkdown h1, .stMarkdown h2 {
        color: #1e3a8a;
    }
    .preview-container {
        border: 2px solid #e0e0e0;
        border-radius: 8px;
        padding: 10px;
        background-color: #ffffff;
    }
    </style>
    """,
    unsafe_allow_html=True
)

# Streamlit app title
st.title("📊 PPT Generator from Outline")

# Initialize session state
if "slides" not in st.session_state:
    st.session_state.slides = []
if "undo_stack" not in st.session_state:
    st.session_state.undo_stack = []
if "redo_stack" not in st.session_state:
    st.session_state.redo_stack = []
if "edit_index" not in st.session_state:
    st.session_state.edit_index = None
if "dataset" not in st.session_state:
    st.session_state.dataset = None
if "sorted_slides" not in st.session_state:
    st.session_state.sorted_slides = []

# Helper function to add text to a shape with formatting
def add_text_to_shape(shape, text, font_name, font_size=18, bold=False, italic=False, font_color=(0, 0, 0), alignment="left"):
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = bold
    font.italic = italic
    font.color.rgb = RGBColor(*font_color)
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[alignment]

# Helper function to add a bullet list with formatting
def add_bullet_list(slide, left, top, width, height, bullets, font_name, font_size=18, bold=False, italic=False, font_color=(0, 0, 0), alignment="left"):
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    for bullet in bullets:
        p = text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = RGBColor(*font_color)
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[alignment]

# Helper function to set slide background
def set_slide_background(slide, bg_type, color1, color2=None):
    background = slide.background
    fill = background.fill
    if bg_type == "Solid":
        fill.solid()
        fill.fore_color.rgb = RGBColor(*color1)
    elif bg_type == "Gradient":
        fill.gradient()
        fill.gradient_angle = 90
        stop1 = fill.gradient_stops[0]
        stop1.color.rgb = RGBColor(*color1)
        stop2 = fill.gradient_stops[1]
        stop2.color.rgb = RGBColor(*color2 if color2 else color1)

# Helper function to generate chart preview (for manual chart input)
def generate_chart_preview(chart_type, categories, values, font_color_hex):
    if not categories or not values:
        return None
    plt.figure(figsize=(4, 3))
    if chart_type == "pie":
        plt.pie(values, labels=categories, autopct='%1.1f%%', colors=['#4CAF50', '#FF9800', '#2196F3'])
    elif chart_type == "bar":
        plt.bar(categories, values, color='#4CAF50')
    elif chart_type == "line":
        plt.plot(categories, values, marker='o', color='#4CAF50')
    plt.title("Chart Preview", color=font_color_hex)
    buf = io.BytesIO()
    plt.savefig(buf, format="png", bbox_inches="tight")
    plt.close()
    buf.seek(0)
    return buf

# Helper function to regenerate Plotly figure from chart data
def regenerate_plotly_fig(df, chart_type, x_col, y_col, category_col, title, font_name, font_color_hex):
    try:
        if chart_type == "bar":
            fig = px.bar(df, x=x_col, y=y_col, color=category_col if category_col else None)
        elif chart_type == "line":
            fig = px.line(df, x=x_col, y=y_col, color=category_col if category_col else None)
        elif chart_type == "pie":
            fig = px.pie(df, names=x_col, values=y_col)
        elif chart_type == "scatter":
            fig = px.scatter(df, x=x_col, y=y_col, color=category_col if category_col else None)
        fig.update_layout(
            title=title or "Chart",
            font=dict(family=font_name, color=font_color_hex)
        )
        return fig
    except Exception as e:
        st.warning(f"⚠️ Failed to regenerate chart: {str(e)}")
        return None

# Helper function to determine slide layout
def determine_layout(slide_data):
    has_title = bool(slide_data.get("title"))
    has_content = bool(slide_data.get("content") or slide_data.get("chart") or slide_data.get("image"))
    if has_title and not has_content:
        return "Title Slide"
    elif has_title and has_content:
        return "Title and Content"
    else:
        return "Blank"

# Tabs for navigation
tab1, tab2, tab3, tab4 = st.tabs(["📝 Add Slides", "🎨 Style Options", "📋 JSON Input", "👀 Preview"])

with tab1:
    # Form for adding/editing slides
    st.header("Add or Edit Slide")
    edit_mode = st.session_state.edit_index is not None
    slide_data = st.session_state.slides[st.session_state.edit_index] if edit_mode else {}
    
    with st.form("slide_form", clear_on_submit=True):
        title = st.text_input(
            "Slide Title", 
            value=slide_data.get("title", ""), 
            placeholder="e.g., Introduction",
            help="Enter the title for your slide."
        )
        # Text formatting controls
        st.subheader("Text Formatting")
        title_font_size = st.selectbox(
            "Title Font Size",
            [20, 24, 28, 32],
            index=[20, 24, 28, 32].index(slide_data.get("title_font_size", 24)) if slide_data.get("title_font_size") else 1,
            help="Font size for the slide title."
        )
        title_alignment = st.selectbox(
            "Title Alignment",
            ["Left", "Center", "Right"],
            index=["Left", "Center", "Right"].index(slide_data.get("title_alignment", "Left")) if slide_data.get("title_alignment") else 0,
            help="Alignment for the slide title."
        )
        title_bold = st.checkbox("Title Bold", value=slide_data.get("title_bold", True), help="Make the title bold.")
        title_italic = st.checkbox("Title Italic", value=slide_data.get("title_italic", False), help="Make the title italic.")
        body_font_size = st.selectbox(
            "Body Font Size",
            [14, 16, 18, 20],
            index=[14, 16, 18, 20].index(slide_data.get("body_font_size", 18)) if slide_data.get("body_font_size") else 2,
            help="Font size for bullet points."
        )
        body_alignment = st.selectbox(
            "Body Alignment",
            ["Left", "Center", "Right"],
            index=["Left", "Center", "Right"].index(slide_data.get("body_alignment", "Left")) if slide_data.get("body_alignment") else 0,
            help="Alignment for bullet points."
        )
        body_bold = st.checkbox("Body Bold", value=slide_data.get("body_bold", False), help="Make bullet points bold.")
        body_italic = st.checkbox("Body Italic", value=slide_data.get("body_italic", False), help="Make bullet points italic.")
        
        bullets = st.multiselect(
            "Bullet Points (add one at a time)",
            options=slide_data.get("content", []) + [""] if edit_mode else [],
            default=slide_data.get("content", []),
            placeholder="Type and select bullet points",
            help="Add bullet points one by one."
        )
        new_bullet = st.text_input("New Bullet Point", placeholder="Type new bullet and add to list")
        if new_bullet and new_bullet not in bullets:
            bullets.append(new_bullet)
        
        # Dataset-driven chart input
        st.subheader("Chart Options")
        chart_input_type = st.radio("Chart Input Method", ["Manual", "Dataset"], index=0)
        
        chart_type = None
        chart_data = None
        plotly_fig = None
        
        if chart_input_type == "Manual":
            chart_type = st.selectbox(
                "Chart Type", 
                ["None", "Pie", "Bar", "Line"],
                index=["None", "Pie", "Bar", "Line"].index(slide_data.get("chart", "None").capitalize()) if slide_data.get("chart") else 0,
                help="Select a chart type or None."
            )
            categories = []
            values = []
            if chart_type != "None":
                st.subheader("Chart Data")
                num_points = st.slider("Number of Data Points", 2, 10, 3)
                cols = st.columns(2)
                with cols[0]:
                    for i in range(num_points):
                        cat = st.text_input(
                            f"Category {i+1}", 
                            value=slide_data.get("chart_data", {}).get("categories", [])[i] if i < len(slide_data.get("chart_data", {}).get("categories", [])) else "",
                            key=f"cat_{i}"
                        )
                        categories.append(cat)
                with cols[1]:
                    for i in range(num_points):
                        val = st.number_input(
                            f"Value {i+1}", 
                            value=slide_data.get("chart_data", {}).get("values", [])[i] if i < len(slide_data.get("chart_data", {}).get("values", [])) else 0.0,
                            key=f"val_{i}"
                        )
                        values.append(val)
                chart_data = {"categories": [c for c in categories if c], "values": [v for v in values if v]}
        
        elif chart_input_type == "Dataset":
            uploaded_file = st.file_uploader("Upload Dataset (CSV)", type=["csv"], key="dataset_upload")
            if uploaded_file:
                try:
                    df = pd.read_csv(uploaded_file)
                    st.session_state.dataset = df
                    st.dataframe(df.head(), use_container_width=True)
                except Exception as e:
                    st.error(f"❌ Invalid CSV file: {str(e)}")
            
            if st.session_state.dataset is not None:
                df = st.session_state.dataset
                chart_type = st.selectbox(
                    "Chart Type",
                    ["Bar", "Line", "Pie", "Scatter"],
                    index=["Bar", "Line", "Pie", "Scatter"].index(slide_data.get("chart", "Bar")) if slide_data.get("chart") else 0,
                    help="Select a chart type for the dataset."
                )
                x_col = st.selectbox("X-Axis Column", df.columns, key="x_col")
                y_col = st.selectbox(
                    "Y-Axis Column",
                    df.select_dtypes(include=['float64', 'int64']).columns,
                    key="y_col"
                )
                category_col = st.selectbox("Category (Optional)", ["None"] + list(df.columns), key="category_col")
                
                # Validate selections
                if not pd.api.types.is_numeric_dtype(df[y_col]):
                    st.error("❌ Y-Axis column must be numeric.")
                elif df[x_col].isna().any() or df[y_col].isna().any():
                    st.error("❌ Selected columns contain missing values.")
                else:
                    # Generate Plotly chart
                    try:
                        if chart_type == "Bar":
                            fig = px.bar(df, x=x_col, y=y_col, color=category_col if category_col != "None" else None)
                        elif chart_type == "Line":
                            fig = px.line(df, x=x_col, y=y_col, color=category_col if category_col != "None" else None)
                        elif chart_type == "Pie":
                            fig = px.pie(df, names=x_col, values=y_col)
                        elif chart_type == "Scatter":
                            fig = px.scatter(df, x=x_col, y=y_col, color=category_col if category_col != "None" else None)
                        fig.update_layout(
                            title=title or "Chart",
                            font=dict(family=st.session_state.get("style", {}).get("font", "Arial"), color=st.session_state.get("style", {}).get("font_color", "#000000"))
                        )
                        st.plotly_chart(fig, use_container_width=True)
                        chart_data = {"x_col": x_col, "y_col": y_col, "category_col": category_col if category_col != "None" else None}
                    except Exception as e:
                        st.error(f"❌ Error generating chart: {str(e)}")
        
        image = st.text_input(
            "Image Filename (from uploaded images)",
            value=slide_data.get("image", ""),
            placeholder="e.g., image1.png",
            help="Enter the filename of an uploaded image."
        )
        slide_transition = st.selectbox(
            "Transition",
            ["None", "Fade", "Push", "Wipe", "Morph", "Zoom"],
            index=["None", "Fade", "Push", "Wipe", "Morph", "Zoom"].index(slide_data.get("transition", "Fade")) if slide_data.get("transition") else 1,
            help="Select a transition effect for the slide."
        )
        submit = st.form_submit_button("Save Slide" if edit_mode else "Add Slide")
        
        if submit:
            if not title.strip():
                st.error("❌ Please provide a slide title.")
            elif chart_input_type == "Dataset" and st.session_state.dataset is None:
                st.error("❌ Please upload a dataset for dataset-driven charts.")
            else:
                new_slide = {
                    "title": title,
                    "content": [b for b in bullets if b],
                    "chart": chart_type.lower() if chart_type and chart_type != "None" else "",
                    "chart_data": chart_data,
                    "chart_input_type": chart_input_type,
                    "image": image.strip() if image.strip() else None,
                    "transition": slide_transition,
                    "title_font_size": title_font_size,
                    "title_alignment": title_alignment.lower(),
                    "title_bold": title_bold,
                    "title_italic": title_italic,
                    "body_font_size": body_font_size,
                    "body_alignment": body_alignment.lower(),
                    "body_bold": body_bold,
                    "body_italic": body_italic
                }
                st.session_state.undo_stack.append(st.session_state.slides.copy())
                if edit_mode:
                    st.session_state.slides[st.session_state.edit_index] = new_slide
                    st.session_state.edit_index = None
                else:
                    st.session_state.slides.append(new_slide)
                st.session_state.redo_stack = []
                st.success("✅ Slide saved!" if edit_mode else "✅ Slide added!")

with tab2:
    # Style options with presets
    st.header("Style Options")
    style_presets = {
        "Professional": {"font": "Calibri", "font_color": "#000080", "bg_type": "Gradient", "bg_color1": "#DDE4FF", "bg_color2": "#FFFFFF"},
        "Creative": {"font": "Roboto", "font_color": "#D81B60", "bg_type": "Solid", "bg_color1": "#FFE082", "bg_color2": None},
        "Minimalist": {"font": "Arial", "font_color": "#000000", "bg_type": "Solid", "bg_color1": "#FFFFFF", "bg_color2": None}
    }
    preset = st.selectbox("Select Style Preset", ["Custom"] + list(style_presets.keys()))
    if preset != "Custom" and st.button("Apply Preset"):
        st.session_state.style = style_presets[preset]
    
    font_name = st.selectbox(
        "Font Type",
        [
            "Arial", "Calibri", "Times New Roman", "Helvetica", "Verdana", "Georgia",
            "Roboto", "Open Sans", "Trebuchet MS", "Tahoma", "Segoe UI", "Palatino Linotype",
            "Garamond", "Book Antiqua", "Courier New", "Consolas", "Impact", "Comic Sans MS"
        ],
        index=0 if "style" not in st.session_state else ["Arial", "Calibri", "Times New Roman", "Helvetica", "Verdana", "Georgia", "Roboto", "Open Sans", "Trebuchet MS", "Tahoma", "Segoe UI", "Palatino Linotype", "Garamond", "Book Antiqua", "Courier New", "Consolas", "Impact", "Comic Sans MS"].index(st.session_state.style.get("font", "Arial")),
        help="Choose a font for your presentation."
    )
    st.info("ℹ️ Use widely available fonts (e.g., Arial, Calibri) for compatibility. Embed custom fonts in PowerPoint via File > Options > Save.")
    font_color_hex = st.color_picker(
        "Font Color", 
        value=st.session_state.style.get("font_color", "#000000") if "style" in st.session_state else "#000000"
    )
    font_color = tuple(int(font_color_hex[i:i+2], 16) for i in (1, 3, 5))
    
    bg_type = st.selectbox(
        "Background Type", 
        ["Solid", "Gradient"],
        index=["Solid", "Gradient"].index(st.session_state.style.get("bg_type", "Solid")) if "style" in st.session_state else 0
    )
    bg_color1_hex = st.color_picker(
        "Background Color 1", 
        value=st.session_state.style.get("bg_color1", "#FFFFFF") if "style" in st.session_state else "#FFFFFF"
    )
    bg_color1 = tuple(int(bg_color1_hex[i:i+2], 16) for i in (1, 3, 5))
    bg_color2_hex = st.color_picker(
        "Background Color 2 (Gradient)", 
        value=st.session_state.style.get("bg_color2", "#DDE4FF") if "style" in st.session_state and st.session_state.style.get("bg_color2") else "#DDE4FF"
    ) if bg_type == "Gradient" else None
    bg_color2 = tuple(int(bg_color2_hex[i:i+2], 16) for i in (1, 3, 5)) if bg_color2_hex else None
    
    # Set default style
    if "style" not in st.session_state:
        st.session_state.style = {
            "font": font_name,
            "font_color": font_color_hex,
            "bg_type": bg_type,
            "bg_color1": bg_color1_hex,
            "bg_color2": bg_color2_hex if bg_color2_hex else None
        }
    
    layout_name = st.selectbox("Default Slide Layout", ["Title Slide", "Title and Content", "Blank"], index=1)
    layout_indices = {"Title Slide": 0, "Title and Content": 1, "Blank": 6}
    default_layout_index = layout_indices[layout_name]
    
    transition = st.selectbox("Default Transition", ["None", "Fade", "Push", "Wipe", "Morph", "Zoom"], index=1)
    
    # Theme export/import
    st.subheader("Export/Import Theme")
    if st.button("Export Theme"):
        theme_json = json.dumps(st.session_state.style, indent=2)
        st.download_button(
            label="Download Theme JSON",
            data=theme_json,
            file_name="theme.json",
            mime="application/json"
        )
    theme_file = st.file_uploader("Upload Theme JSON", type=["json"])
    if theme_file:
        try:
            theme_data = json.load(theme_file)
            required_keys = ["font", "font_color", "bg_type", "bg_color1"]
            if all(key in theme_data for key in required_keys):
                st.session_state.style = theme_data
                st.success("✅ Theme imported!")
            else:
                st.error("❌ Invalid theme JSON: Missing required keys.")
        except json.JSONDecodeError:
            st.error("❌ Invalid JSON format.")
    
    st.header("Style Preview")
    st.markdown(
        f"""
        <div style='font-family:{font_name}; font-size:{body_font_size}px; color:{font_color_hex}; background-color:{bg_color1_hex}; padding:10px; border-radius:8px;'>
            Sample Text (Font: {font_name}, Size: {body_font_size}pt, Color: {font_color_hex})
        </div>
        """,
        unsafe_allow_html=True
    )

with tab3:
    # JSON input
    st.header("Paste JSON Outline (Advanced)")
    st.write("Paste a JSON outline to replace current slides. Example:")
    st.code(
        '''
[
  {
    "title": "Introduction",
    "content": ["Point 1", "Point 2"],
    "chart": "pie",
    "chart_data": {"categories": ["Category A", "Category B"], "values": [60, 40]},
    "chart_input_type": "Manual",
    "image": "image1.png",
    "transition": "Fade",
    "title_font_size": 24,
    "title_alignment": "left",
    "title_bold": true,
    "title_italic": false,
    "body_font_size": 18,
    "body_alignment": "left",
    "body_bold": false,
    "body_italic": false
  },
  {
    "title": "Data Analysis",
    "content": ["Analysis Results"],
    "chart": "bar",
    "chart_data": {"x_col": "Region", "y_col": "Sales", "category_col": "Year"},
    "chart_input_type": "Dataset",
    "transition": "Wipe",
    "title_font_size": 28,
    "title_alignment": "center",
    "title_bold": true,
    "title_italic": false,
    "body_font_size": 16,
    "body_alignment": "left",
    "body_bold": false,
    "body_italic": false
  }
]
        ''',
        language="json"
    )
    outline_json = st.text_area("Enter PPT Outline (JSON)", height=200, placeholder="Paste your JSON outline here")
    if st.button("Load JSON Outline"):
        if not outline_json.strip():
            st.error("❌ Please provide a JSON outline.")
        else:
            try:
                slides = json.loads(outline_json)
                if not isinstance(slides, list):
                    st.error("❌ Outline must be a list of slides.")
                else:
                    st.session_state.undo_stack.append(st.session_state.slides.copy())
                    st.session_state.slides = slides
                    st.session_state.redo_stack = []
                    st.success("✅ JSON outline loaded!")
            except json.JSONDecodeError:
                st.error("❌ Invalid JSON format.")

with tab4:
    # Slide preview and management
    st.header("Preview and Manage Slides")
    uploaded_images = st.file_uploader("Upload Images (Optional)", type=["png", "jpg"], accept_multiple_files=True, help="Upload images to use in slides.")
    image_files = {f.name: f for f in (uploaded_images or [])}
    
    if st.session_state.slides:
        st.subheader("Current Slides")
        # Drag-and-drop reordering
        slide_items = [{"id": i, "title": slide.get("title", "Untitled")} for i, slide in enumerate(st.session_state.slides)]
        sorted_items = sort_items(slide_items, multi_containers=False, direction="vertical", key="slide_sorter")
        if sorted_items != slide_items:
            new_order = [item["id"] for item in sorted_items]
            st.session_state.undo_stack.append(st.session_state.slides.copy())
            st.session_state.slides = [st.session_state.slides[i] for i in new_order]
            st.session_state.redo_stack = []
            st.success("✅ Slides reordered!")
        
        for i, slide in enumerate(st.session_state.slides):
            with st.expander(f"Slide {i+1}: {slide.get('title', 'Untitled')}"):
                st.markdown(f"**Title**: {slide.get('title', 'Untitled')}")
                if slide.get("content"):
                    st.markdown("**Bullet Points**:")
                    for bullet in slide.get("content", []):
                        st.markdown(f"- {bullet}")
                if slide.get("chart") and slide.get("chart_data"):
                    if slide.get("chart_input_type") == "Manual":
                        chart_buf = generate_chart_preview(
                            slide["chart"], 
                            slide["chart_data"]["categories"], 
                            slide["chart_data"]["values"], 
                            st.session_state.style.get("font_color", "#000000")
                        )
                        if chart_buf:
                            st.image(chart_buf, caption="Chart Preview", use_container_width=True)
                    elif slide.get("chart_input_type") == "Dataset" and st.session_state.dataset is not None:
                        df = st.session_state.dataset
                        chart_data = slide.get("chart_data", {})
                        x_col = chart_data.get("x_col")
                        y_col = chart_data.get("y_col")
                        category_col = chart_data.get("category_col")
                        if x_col and y_col and x_col in df.columns and y_col in df.columns:
                            fig = regenerate_plotly_fig(
                                df, 
                                slide["chart"], 
                                x_col, 
                                y_col, 
                                category_col, 
                                slide["title"], 
                                st.session_state.style.get("font", "Arial"), 
                                st.session_state.style.get("font_color", "#000000")
                            )
                            if fig:
                                st.plotly_chart(fig, use_container_width=True)
                        else:
                            st.warning("⚠️ Invalid chart data or dataset missing.")
                if slide.get("image") and slide.get("image") in image_files:
                    st.image(image_files[slide["image"]], caption="Image Preview", width=200)
                st.markdown(f"**Transition**: {slide.get('transition', transition)}")
                col1, col2, col3 = st.columns(3)
                with col1:
                    if st.button("Edit", key=f"edit_{i}"):
                        st.session_state.edit_index = i
                        st.experimental_rerun()
                with col2:
                    if st.button("Delete", key=f"delete_{i}"):
                        st.session_state.undo_stack.append(st.session_state.slides.copy())
                        st.session_state.slides.pop(i)
                        st.session_state.redo_stack = []
                        st.success("✅ Slide deleted!")
                with col3:
                    if st.button("Duplicate", key=f"duplicate_{i}"):
                        st.session_state.undo_stack.append(st.session_state.slides.copy())
                        st.session_state.slides.append(slide.copy())
                        st.session_state.redo_stack = []
                        st.success("✅ Slide duplicated!")
    
        # Undo/Redo buttons
        col1, col2 = st.columns(2)
        with col1:
            if st.button("Undo", disabled=not st.session_state.undo_stack):
                st.session_state.redo_stack.append(st.session_state.slides.copy())
                st.session_state.slides = st.session_state.undo_stack.pop()
                st.success("✅ Undone!")
        with col2:
            if st.button("Redo", disabled=not st.session_state.redo_stack):
                st.session_state.undo_stack.append(st.session_state.slides.copy())
                st.session_state.slides = st.session_state.redo_stack.pop()
                st.success("✅ Redone!")
    
    else:
        st.info("No slides added yet. Use the 'Add Slides' tab to create slides.")

# Generate PPT
if st.button("Generate PPT", key="generate_ppt"):
    if not st.session_state.slides:
        st.error("❌ Please add at least one slide.")
    else:
        try:
            prs = Presentation()
            for slide_data in st.session_state.slides:
                if not isinstance(slide_data, dict):
                    st.warning(f"⚠️ Skipping invalid slide data: {slide_data}")
                    continue
                
                title = slide_data.get("title", "Untitled")
                content = slide_data.get("content", [])
                chart_type = slide_data.get("chart", "").lower()
                chart_data_input = slide_data.get("chart_data", None)
                chart_input_type = slide_data.get("chart_input_type", "Manual")
                image_path = slide_data.get("image", None)
                slide_transition = slide_data.get("transition", transition)
                title_font_size = slide_data.get("title_font_size", 24)
                title_alignment = slide_data.get("title_alignment", "left")
                title_bold = slide_data.get("title_bold", True)
                title_italic = slide_data.get("title_italic", False)
                body_font_size = slide_data.get("body_font_size", 18)
                body_alignment = slide_data.get("body_alignment", "left")
                body_bold = slide_data.get("body_bold", False)
                body_italic = slide_data.get("body_italic", False)
                
                # Auto-layout detection
                layout_name = determine_layout(slide_data)
                layout_index = layout_indices[layout_name]
                slide_layout = prs.slide_layouts[layout_index]
                slide = prs.slides.add_slide(slide_layout)
                set_slide_background(slide, bg_type, bg_color1, bg_color2)
                
                if slide.shapes.title:
                    title_shape = slide.shapes.title
                else:
                    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                add_text_to_shape(
                    title_shape, title, font_name=font_name, font_size=title_font_size,
                    bold=title_bold, italic=title_italic, font_color=font_color, alignment=title_alignment
                )
                
                if content:
                    add_bullet_list(
                        slide, Inches(0.5), Inches(1.5), Inches(4), Inches(4), content,
                        font_name=font_name, font_size=body_font_size, bold=body_bold,
                        italic=body_italic, font_color=font_color, alignment=body_alignment
                    )
                
                if chart_type in ["pie", "bar", "line"] and chart_data_input and chart_input_type == "Manual":
                    try:
                        chart_data = CategoryChartData()
                        chart_data.categories = chart_data_input["categories"]
                        chart_data.add_series("Data", chart_data_input["values"])
                        chart_type_enum = {
                            "pie": XL_CHART_TYPE.PIE,
                            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
                            "line": XL_CHART_TYPE.LINE
                        }[chart_type]
                        chart = slide.shapes.add_chart(
                            chart_type_enum, Inches(5), Inches(1.5), Inches(4), Inches(3), chart_data
                        ).chart
                        chart.has_title = True
                        chart.chart_title.text_frame.text = f"{title} Chart"
                        p = chart.chart_title.text_frame.paragraphs[0]
                        p.font.name = font_name
                        p.font.size = Pt(14)
                        p.font.color.rgb = RGBColor(*font_color)
                    except (KeyError, TypeError, ValueError) as e:
                        st.warning(f"⚠️ Invalid chart data for slide '{title}': {str(e)}")
                
                if chart_type in ["bar", "line", "pie", "scatter"] and chart_input_type == "Dataset" and chart_data_input and st.session_state.dataset is not None:
                    try:
                        df = st.session_state.dataset
                        x_col = chart_data_input.get("x_col")
                        y_col = chart_data_input.get("y_col")
                        category_col = chart_data_input.get("category_col")
                        if x_col and y_col and x_col in df.columns and y_col in df.columns:
                            fig = regenerate_plotly_fig(
                                df, chart_type, x_col, y_col, category_col, title,
                                st.session_state.style.get("font", "Arial"),
                                st.session_state.style.get("font_color", "#000000")
                            )
                            if fig:
                                chart_stream = io.BytesIO()
                                fig.write_image(file=chart_stream, format="png")
                                chart_stream.seek(0)
                                slide.shapes.add_picture(chart_stream, Inches(5), Inches(1.5), width=Inches(4))
                            else:
                                st.warning(f"⚠️ Failed to generate chart for slide '{title}'.")
                        else:
                            st.warning(f"⚠️ Invalid chart data for slide '{title}'.")
                    except Exception as e:
                        st.warning(f"⚠️ Failed to add dataset chart for slide '{title}': {str(e)}")
                
                if image_path and image_path in image_files:
                    try:
                        img_stream = io.BytesIO(image_files[image_path].read())
                        slide.shapes.add_picture(img_stream, Inches(0.5), Inches(3.5), width=Inches(4))
                    except Exception as e:
                        st.warning(f"⚠️ Failed to add image '{image_path}': {str(e)}")
                
                slide.notes_slide.notes_text_frame.text = f"Recommended transition: {slide_transition}"
            
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            
            st.download_button(
                label="📥 Download PPT",
                data=buffer,
                file_name="Generated_Presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            st.success("✅ PPT generated successfully!")
        
        except ValueError as e:
            st.error(f"❌ Error in slide structure: {str(e)}")
        except Exception as e:
            st.error(f"❌ Error generating PPT: {str(e)}")

# Tutorial button
if st.button("Show Tutorial"):
    st.markdown("""
    **Welcome to the PPT Generator!** Follow these steps:
    1. **Add Slides**: Create slides with titles, bullet points, charts (manual or dataset-driven), and images. Customize text formatting per slide.
    2. **Customize Styles**: Choose fonts, colors, and backgrounds. Export/import themes as JSON.
    3. **Advanced Input**: Paste a JSON outline in the 'JSON Input' tab.
    4. **Preview and Edit**: View, edit, delete, duplicate, or reorder slides (drag-and-drop) in the 'Preview' tab.
    5. **Generate PPT**: Click 'Generate PPT' to download your presentation.
    """)
